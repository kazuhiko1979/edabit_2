# Location名はpres2が原因かどうか検証　20240607
import os
import zipfile
import shutil
import re
from spire.presentation import Presentation, FileFormat
import ctypes
from pptx import Presentation as pptx_Presentation  # pptxライブラリをインポート
from pptx.util import Inches, Pt  # インチとポイントを扱うためのユーティリティをインポート
from pptx.dml.color import RGBColor


# 変数に格納先のパスを設定
ZIP_FILE_DIR = "C:\\"

# 対象ファイル名を取得、正規表現を修正
zip_file_name_pattern = re.compile(r"webex_reports_(\d{6})_\d+\.zip")
zip_file_name = None
for file in os.listdir(ZIP_FILE_DIR):
    match = zip_file_name_pattern.match(file)
    if match:
        zip_file_name = file
        target_period_pass = match.group(1)
        break

if not zip_file_name:
    print("対象ファイルが見つかりませんでした")
    exit()

zip_file_path = os.path.join(ZIP_FILE_DIR, zip_file_name)

if not os.path.exists(zip_file_path):
    print("zipファイルが存在しません")
    exit()

# パスワードの設定
password = f"Managed-lan-{target_period_pass}"
password_bytes = password.encode()

# 出力パスと一時ディレクトリの設定
OUTPUT_PATH = os.path.join(ZIP_FILE_DIR, f"output_{target_period_pass}")
os.makedirs(OUTPUT_PATH, exist_ok=True)

temp_dir = os.path.join(OUTPUT_PATH, "all")
os.makedirs(temp_dir, exist_ok=True)

# 圧縮ファイルを一時ディレクトリに展開
with zipfile.ZipFile(zip_file_path, 'r') as zip_ref:
    zip_ref.extractall(temp_dir, pwd=password_bytes)

# 各 ORG フォルダを元にクライアントディレクトリを作成し、ファイルを移動する
for org_name in os.listdir(temp_dir):
    org_path = os.path.join(temp_dir, org_name)
    if os.path.isdir(org_path):
        client_dir = os.path.join(OUTPUT_PATH, org_name)
        os.makedirs(client_dir, exist_ok=True)

        for file_name in os.listdir(org_path):
            if file_name.endswith(('.pptx', '.ppt')):
                src_path = os.path.join(org_path, file_name)

                # ファイル名の先頭部分（最初の「-」の手前の文字列）を取得
                base_name = file_name.split('-')[0]
                if base_name == org_name:
                    dst_path = os.path.join(client_dir, file_name)
                else:
                    # フォルダ名とファイルの先頭部分が異なる場合、対応するフォルダに移動
                    other_client_dir = os.path.join(OUTPUT_PATH, base_name)
                    os.makedirs(other_client_dir, exist_ok=True)
                    dst_path = os.path.join(other_client_dir, file_name)
                shutil.move(src_path, dst_path)

        # 元のフォルダを削除
        shutil.rmtree(client_dir)

# 一時ディレクトリを削除
shutil.rmtree(temp_dir)


# 各クライアントディレクトリ内で処理を実行
for client_dir in [d for d in os.listdir(OUTPUT_PATH) if os.path.isdir(os.path.join(OUTPUT_PATH, d))]:
    client_path = os.path.join(OUTPUT_PATH, client_dir)
    files = [f for f in os.listdir(client_path) if f.endswith(('.pptx', '.ppt'))]

    if len(files) < 2:
        print(f"2つ以上の PowerPoint ファイルが必要です。'{client_dir}' にあるファイル数: {len(files)}")
        continue

    # ORG名とLocationを抽出
    org_name = client_dir.split('-')[0]
    locations = [f.split('_')[0] for f in files]

    first_file_path = os.path.join(client_path, files[0])
    pres1 = Presentation()
    pres1.LoadFromFile(first_file_path)
    master_slide_7 = pres1.Slides.get_Item(6)
    pres1.Slides.RemoveAt(6)

    for i, filename in enumerate(files[1:], 1):
        file_path = os.path.join(client_path, filename)
        pres2 = Presentation()
        pres2.LoadFromFile(file_path)

        for j in range(2, 6):
            slide = pres2.Slides.get_Item(j)
            pres1.Slides.AppendBySlide(slide)

            # スライドにLocationをテキストで挿入
            pptx_slide = pptx_Presentation(file_path).slides[j - 1]  # pptxで読み込む
            text_box = pptx_slide.shapes.add_textbox(Inches(4), Inches(0.5), Inches(8), Inches(1))  # テキストボックスを追加
            text_frame = text_box.text_frame
            text_frame.text = locations[i - 1]  # Locationをテキストとして設定
            text_frame.paragraphs[0].font.size = Pt(30)  # フォントサイズを30ptに設定
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # テキストの色を白に設定

            print(text_frame.text)

        # pres2を一時ファイルとして保存
        temp_file_path = os.path.join(client_path, f"temp_{filename}")
        pres2.SaveToFile(temp_file_path, FileFormat.Pptx2019)
        pres2.Dispose()

        # 一時ファイルから元のファイルを上書き
        os.remove(file_path)
        os.rename(temp_file_path, file_path)

    pres1.Slides.AppendBySlide(master_slide_7)
    output_path = os.path.join(client_path, f"{client_dir}.pptx")
    pres1.SaveToFile(output_path, FileFormat.Pptx2019)

    if os.path.getsize(output_path) > 18 * 1024 * 1024:
        part_num = 1
        with open(output_path, 'rb') as f:
            while chunk := f.read(18 * 1024 * 1024):
                part_name = f"{client_dir}_{part_num}.pptx"
                part_path = os.path.join(client_path, part_name)

                with open(part_path, 'wb') as part_file:
                    part_file.write(chunk)
                part_num += 1

        os.remove(output_path)
    else:
        pres1.SaveToFile(output_path, FileFormat.Pptx2019)

    pres1.Dispose()

    for filename in files:
        file_path = os.path.join(client_path, filename)
        os.remove(file_path)

    print(f"'{client_dir}' の処理が完了しました。")

    # 対象ファイルの属性を変更して削除
    if os.path.exists(zip_file_path):
        # ファイルの読み取り専用属性を解除
        FILE_ATTRIBUTE_NORMAL = 0x80
        ctypes.windll.kernel32.SetFileAttributesW(zip_file_path, FILE_ATTRIBUTE_NORMAL)

        try:
            os.remove(zip_file_path)
            print(f"対象ファイル '{zip_file_name}' が削除されました。")
        except PermissionError as e:
            print(f"ファイル削除中にエラーが発生しました: {e}")

    print("すべての処理が完了しました。")



### 2024/06/06 Location名挿入処理前オリジナル　####

# import os
# import zipfile
# import shutil
# import re
# from spire.presentation import Presentation, FileFormat
# import ctypes
#
#
# # 必要なライブラリをインストール: RW-FAT環境実行時
# # proxy = "--proxy=10.160.120.156:8080"
# # os.system(f"pip install {proxy} spire-presentation")
# # zipfile shutil は環境によってインストール
# # os.system(f"pip install {proxy} zipfile shutil spire-presentation")
#
# # 変数に格納先のパスを設定
# ZIP_FILE_DIR = "C:\\"
#
#
# # 対象ファイル名を取得、正規表現を修正
# zip_file_name_pattern = re.compile(r"webex_reports_(\d{6})_\d+\.zip")
# zip_file_name = None
# for file in os.listdir(ZIP_FILE_DIR):
#     match = zip_file_name_pattern.match(file)
#     if match:
#         zip_file_name = file
#         target_period_pass = match.group(1)
#         break
#
# if not zip_file_name:
#     print("対象ファイルが見つかりませんでした")
#     exit()
#
# zip_file_path = os.path.join(ZIP_FILE_DIR, zip_file_name)
#
# if not os.path.exists(zip_file_path):
#     print("zipファイルが存在しません")
#     exit()
#
# # パスワードの設定
# password = f"Managed-lan-{target_period_pass}"
# password_bytes = password.encode()
#
# # 出力パスと一時ディレクトリの設定
# OUTPUT_PATH = os.path.join(ZIP_FILE_DIR, f"output_{target_period_pass}")
# os.makedirs(OUTPUT_PATH, exist_ok=True)
#
# temp_dir = os.path.join(OUTPUT_PATH, "all")
# os.makedirs(temp_dir, exist_ok=True)
#
# # 圧縮ファイルを一時ディレクトリに展開
# with zipfile.ZipFile(zip_file_path, 'r') as zip_ref:
#     zip_ref.extractall(temp_dir, pwd=password_bytes)
#
#
# # 各 ORG フォルダを元にクライアントディレクトリを作成し、ファイルを移動する
# for org_name in os.listdir(temp_dir):
#     org_path = os.path.join(temp_dir, org_name)
#     if os.path.isdir(org_path):
#         client_dir = os.path.join(OUTPUT_PATH, org_name)
#         os.makedirs(client_dir, exist_ok=True)
#
#         for file_name in os.listdir(org_path):
#             if file_name.endswith(('.pptx', '.ppt')):
#                 src_path = os.path.join(org_path, file_name)
#
#                 # ファイル名の先頭部分（最初の「-」の手前の文字列）を取得
#                 base_name = file_name.split('-')[0]
#                 if base_name == org_name:
#                     dst_path = os.path.join(client_dir, file_name)
#                 else:
#                     # フォルダ名とファイルの先頭部分が異なる場合、対応するフォルダに移動
#                     other_client_dir = os.path.join(OUTPUT_PATH, base_name)
#                     os.makedirs(other_client_dir, exist_ok=True)
#                     dst_path = os.path.join(other_client_dir, file_name)
#                 shutil.move(src_path, dst_path)
#
#         # 元のフォルダを削除
#         shutil.rmtree(client_dir)
#
# # 一時ディレクトリを削除
# shutil.rmtree(temp_dir)
#
#
# # 各クライアントディレクトリ内で処理を実行
# for client_dir in [d for d in os.listdir(OUTPUT_PATH) if os.path.isdir(os.path.join(OUTPUT_PATH, d))]:
#     client_path = os.path.join(OUTPUT_PATH, client_dir)
#     files = [f for f in os.listdir(client_path) if f.endswith(('.pptx', '.ppt'))]
#
#     if len(files) < 2:
#         print(f"2つ以上の PowerPoint ファイルが必要です。'{client_dir}' にあるファイル数: {len(files)}")
#         continue
#
#     first_file_path = os.path.join(client_path, files[0])
#     pres1 = Presentation()
#     pres1.LoadFromFile(first_file_path)
#     master_slide_7 = pres1.Slides.get_Item(6)
#     pres1.Slides.RemoveAt(6)
#
#     for filename in files[1:]:
#         file_path = os.path.join(client_path, filename)
#         pres2 = Presentation()
#         pres2.LoadFromFile(file_path)
#
#         for i in range(2, 6):
#             slide = pres2.Slides.get_Item(i)
#             pres1.Slides.AppendBySlide(slide)
#         pres2.Dispose()
#
#     pres1.Slides.AppendBySlide(master_slide_7)
#     output_path = os.path.join(client_path, f"{client_dir}.pptx")
#     pres1.SaveToFile(output_path, FileFormat.Pptx2019)
#
#     if os.path.getsize(output_path) > 18 * 1024 * 1024:
#         part_num = 1
#         with open(output_path, 'rb') as f:
#             while chunk := f.read(18 * 1024 * 1024):
#                 part_name = f"{client_dir}_{part_num}.pptx"
#                 part_path = os.path.join(client_path, part_name)
#
#                 with open(part_path, 'wb') as part_file:
#                     part_file.write(chunk)
#                 part_num += 1
#
#         os.remove(output_path)
#     else:
#         pres1.SaveToFile(output_path, FileFormat.Pptx2019)
#
#     pres1.Dispose()
#
#     for filename in files:
#         file_path = os.path.join(client_path, filename)
#         os.remove(file_path)
#
#     print(f"'{client_dir}' の処理が完了しました。")
#
# # 対象ファイルの属性を変更して削除
# if os.path.exists(zip_file_path):
#     # ファイルの読み取り専用属性を解除
#     FILE_ATTRIBUTE_NORMAL = 0x80
#     ctypes.windll.kernel32.SetFileAttributesW(zip_file_path,
# FILE_ATTRIBUTE_NORMAL)
#     try:
#         os.remove(zip_file_path)
#         print(f"対象ファイル '{zip_file_name}' が削除されました。")
#     except PermissionError as e:
#         print(f"ファイル削除中にエラーが発生しました: {e}")
#
#
# print("すべての処理が完了しました。")